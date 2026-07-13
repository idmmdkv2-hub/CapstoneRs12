import streamlit as st
import json
import google.genai as genai
from pptx import Presentation
from docx import Document
from PIL import Image

# Initialize the Gemini Client
client = genai.Client(api_key="AQ.Ab8RN6J7CZxd3T-UCaKQzoFF_1rsD4yZbP8kuE-08WaF2KuDvQ")

# Helper function to extract text from different file types
def extract_text(uploaded_file):
    file_ext = uploaded_file.name.split('.')[-1].lower()
    
    if file_ext == 'docx':
        doc = Document(uploaded_file)
        return "\n".join([p.text for p in doc.paragraphs])
        
    elif file_ext == 'pptx':
        prs = Presentation(uploaded_file)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
        
    elif file_ext in ['png', 'jpg', 'jpeg']:
        # If it's an image, we return the PIL Image object directly so Gemini can scan it visually
        return Image.open(uploaded_file)
        
    return None

# --- STREAMLIT WEB UI CONFIG ---
st.set_page_config(page_title="AI Study Gizmo", page_icon="📚", layout="centered")
st.title("📚 AI Study Scanner")
st.write("Upload a Presentation, Word Doc, or Photo to generate an interactive study questionnaire instantly!")

# File uploader widget
uploaded_file = st.file_uploader(
    "Choose a file", 
    type=["docx", "pptx", "png", "jpg", "jpeg"]
)

if uploaded_file is not None:
    st.success(f"Successfully uploaded: {uploaded_file.name}")
    
    # 1. Extract the text/image content
    content = extract_text(uploaded_file)
    
    # 2. Trigger Generation Button
    if st.button("✨ Generate Questionnaire"):
        with st.spinner("Analyzing materials and crafting questions..."):
            try:
                # Prepare prompt for the AI
                prompt = (
                    "Analyze the attached study material. Generate a study deck of exactly 3 different flashcards based on the core topics. "
                    "Format your entire response strictly as a valid JSON array of objects, containing 'q' (question) and 'a' (answer) keys, like this:\n"
                    '[{"q": "Question 1", "a": "Answer 1"}, {"q": "Question 2", "a": "Answer 2"}]'
                )
                
                # Bundle file contents depending on if it's text string or a direct image asset
                if isinstance(content, Image.Image):
                    # Multimodal call: pass the photo and the prompt together
                    response = client.models.generate_content(
                        model='gemini-2.0-flash',
                        contents=[content, prompt]
                    )
                else:
                    # Text-based call (from docx/pptx)
                    response = client.models.generate_content(
                        model='gemini-2.0-flash',
                        contents=f"Material:\n{content}\n\n{prompt}"
                    )
                
                # Clean up json format strings if returned inside Markdown wrappers
                clean_text = response.text.strip().removeprefix("```json").removesuffix("```").strip()
                
                # Save the questions into the website's temporary memory session
                st.session_state['quiz_data'] = json.loads(clean_text)
                st.session_state['card_index'] = 0
                st.session_state['show_answer'] = False
                
            except Exception as e:
                st.error(f"Failed to build quiz. Server might be busy: {e}")

# --- INTERACTIVE FLASHCARD INTERFACE ---
if 'quiz_data' in st.session_state and st.session_state['quiz_data']:
    quiz = st.session_state['quiz_data']
    idx = st.session_state['card_index']
    
    st.divider()
    st.subheader(f"Question Card {idx + 1} of {len(quiz)}")
    
    # Display the Question Box
    st.info(f"**QUESTION:**\n{quiz[idx]['q']}")
    
    # Reveal Button
    if st.button("👁️ Flip / Reveal Answer"):
        st.session_state['show_answer'] = True
        
    # Display the Answer Box if flipped
    if st.session_state['show_answer']:
        st.success(f"**ANSWER / breakdown:**\n{quiz[idx]['a']}")
        
        # Next Navigation Button
        if idx + 1 < len(quiz):
            if st.button("➡️ Next Question"):
                st.session_state['card_index'] += 1
                st.session_state['show_answer'] = False
                st.rerun()
        else:
            st.balloons()
            st.write("🎉 You completed the deck! Upload another file to start fresh.")